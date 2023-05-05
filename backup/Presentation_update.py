#!/usr/bin/env python
# _*_ coding:utf-8 _*_
import datetime as dt

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

import Template_generator as Tg
from Monthly_parser import monthly_data_arrangement, monthly_parts_and_title, monthly_show_parts, \
    last_two_row_difference, multi_idx_df_multiply_one_col
from WeekTime_parser import short_month
from chart import ChartIOGen

items = ["Q'ty", "Yield", "FSV", "BSV", "CP"]
side_correction = {'FSV': '正檢', 'BSV': '背檢', 'CP': 'CP'}
item_map = {"Yield": 'final yield average', "FSV": 'FSV Loss', "BSV": 'BSV Loss', "CP": 'CP Loss'}
red = RGBColor(255, 0, 0)
green = RGBColor(0, 128, 0)
blue = RGBColor(0, 0, 204)
purple = RGBColor(112, 48, 160)
zero64 = np.float64(0)
percent_10 = np.float64(0.1)


class ActionItemNotFoundError(Exception):
    def __init__(self, part, side):
        self.value = f"Can't find action item for {part}, {side}. Please confirm the content in action_item_for_ppt.xlsx is correct"


def picture_set_up(pic):
    border = pic.line
    border.color.rgb = RGBColor(128, 128, 128)
    border.width = Pt(1.5)


def per_deci(val, operator=False):
    """percentage_decimal_determine"""
    if operator:
        sign = '+'
    else:
        sign = ''
    if val >= percent_10:
        return '{:' + sign + '.1%}'
    else:
        return '{:' + sign + '.2%}'


def monthly_report_gen(date, force_month_summary=False, ignored_parts_and_wafers=None,
                       comparing_to_the_last_year=False, specific_date=False):
    if ignored_parts_and_wafers is not None and specific_date:
        sub_df_dict, wafer_list_dict = ignored_parts_and_wafers
        specific_year = date.year
        specific_month = short_month(date)
    elif specific_date:
        sub_df_dict = None
        wafer_list_dict = None
        specific_year = date.year
        specific_month = short_month(date)
    else:
        sub_df_dict = None
        wafer_list_dict = None
        specific_year = None
        specific_month = None
    monthly_data_dict = monthly_data_arrangement(specific_year=specific_year, specific_month=specific_month,
                                                 sub_df_dict=sub_df_dict, wafer_list_dict=wafer_list_dict,
                                                 compare_year=comparing_to_the_last_year)
    prs = Presentation(r'templates/template.pptx')

    # add_title_slide

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    if force_month_summary:
        show_date = Tg.current_date
        title_day = '-{}'.format(show_date.day)
        day = '/{}'.format(show_date.day)
    else:
        title_day = ''
        day = ''
        show_date = date
    target_month = short_month(show_date)
    title.text = "Yield Status of Conti Product \n– Monthly Summary ({}/{}{})".format(show_date.year, target_month,
                                                                                      title_day)

    p = subtitle.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Product Part:\n'
    font = run.font
    font.bold = True
    font.italic = True
    font.underline = True
    for part in list(monthly_data_dict):
        dfs_last_month = monthly_data_dict[part]['last_month']
        # df = monthly_data_dict[part]['df']
        if dfs_last_month != target_month:
            monthly_data_dict.pop(part, None)  # del part that does't have data yet or is not last month part's data
            continue
        description = monthly_parts_and_title.get(part)
        if description is None:
            continue
        run = p.add_run()
        run.text = part + '(' + description + ')' + '\n'

    for part, subtitle in monthly_show_parts.items():
        part_info_dict = monthly_data_dict.get(part)
        if part_info_dict is None:
            continue

        df = part_info_dict['df']
        chart = ChartIOGen(df)
        monthly_bar_chart = chart.monthly_data()
        last_idx = df.last_valid_index()
        if comparing_to_the_last_year:
            df_last_year = part_info_dict['df_last_year']
            if df_last_year.empty:
                df_compare_difference = df_last_year
            else:
                multi_idx_df_multiply_one_col(df_last_year)
                df_last_year = df_last_year.sum().to_frame('Last_year').T
                multi_idx_df_multiply_one_col(df_last_year, operator='/')
                df_diff = pd.concat([df_last_year, df.tail(1)])

                df_compare_difference = last_two_row_difference(df_diff)
        else:
            df_compare_difference = last_two_row_difference(df)

        # first slide for comprehensive summary

        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Yield trend chart of {} {} ({}/{}{})".format(part, subtitle, show_date.year, show_date.month, day)
        content = slide.placeholders[10]
        tf = content.text_frame
        pcs = df.tail(1).at[last_idx, ('SUMMARY', "Q'ty")]
        tf.text = '{} total output: {:.0f}pcs'.format(part, pcs)
        for item in items:
            if item == "Q'ty":
                continue
            elif item == 'Yield':
                neg_color = red  # r
                pos_color = green  # g
                header = part + ' '
            else:
                neg_color = green
                pos_color = red
                header = ''

            item_val = df.tail(1).loc[last_idx, ('SUMMARY', item)]
            formatted_str = ('{}{}: ' + per_deci(item_val)).format(header, item_map[item], item_val)
            p = tf.add_paragraph()
            p.text = formatted_str
            p.space_before = 0
            if not df_compare_difference.empty:
                run = p.add_run()
                item_diff = df_compare_difference.loc[last_idx, ('SUMMARY', item)]
                run.text = (' (' + per_deci(item_diff, operator=True) + ')').format(item_diff)
                run.font.color.rgb = pos_color if item_diff > zero64 else neg_color
        picture = slide.placeholders[1]
        placeholder_picture = picture.insert_picture(monthly_bar_chart.ref)
        picture_set_up(placeholder_picture)

        # second & third slide of BSV & FSV pareto

        for side in ['BSV', 'FSV']:
            pareto_chart, sorted_data = chart.weekly_defect_data('pareto', side, return_sorted_data=True)

            # force glass align shift item appear in AP1C0 pareto monthly summary
            if part == 'AP1C0' and side == 'BSV':
                if 'Glass Align Shift' not in sorted_data.columns:
                    sorted_data['Glass Align Shift'] = [0]

            pareto_monthly_summary_chart = chart.monthly_data(monthly_pareto_summary=True, side=side,
                                                              title=f'{side} pareto monthly summary',
                                                              sorted_data=sorted_data)
            pareto_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(pareto_slide_layout)
            title = slide.shapes.title

            picture1 = slide.placeholders[1]
            placeholder_picture1 = picture1.insert_picture(pareto_chart.ref)
            picture_set_up(placeholder_picture1)

            picture2 = slide.placeholders[11]
            placeholder_picture2 = picture2.insert_picture(pareto_monthly_summary_chart.ref)
            picture_set_up(placeholder_picture2)

            if '&' in part:
                subtitle = subtitle.replace(" pressure sensor", "")
            title.text = "{} {} Defect Pareto for {} ({}/{}{})".format(part, subtitle, side, show_date.year,
                                                                       show_date.month, day)
            content = slide.placeholders[10]
            tf = content.text_frame
            tf.clear()
            # first line
            p = tf.paragraphs[0]
            side_val = df.tail(1).loc[last_idx, ('SUMMARY', side)]
            p.text = ('Total {} yield loss: ' + per_deci(side_val)).format(side, side_val)
            if not df_compare_difference.empty:
                run = p.add_run()
                side_diff = df_compare_difference.loc[last_idx, ('SUMMARY', side)]
                run.text = (' (' + per_deci(side_diff, operator=True) + ')').format(side_diff)
                run.font.color.rgb = red if side_diff > zero64 else green
            p.space_before = 0

            # p_empty = tf.add_paragraph()
            # p_empty.level = 3
            # p_empty.space_before = 0

            for item in sorted_data.iloc[:, :3]:
                p = tf.add_paragraph()
                item_val = sorted_data.at[0, item]
                p.text = ("{}: " + per_deci(item_val)).format(item, item_val)
                p.space_before = 0
                p.level = 1
                if not df_compare_difference.empty:
                    run = p.add_run()
                    item_diff = df_compare_difference.loc[last_idx, (side, item)]
                    run.text = (' (' + per_deci(item_diff, operator=True) + ')').format(item_diff)
                    run.font.color.rgb = red if item_diff > zero64 else green
                p = tf.add_paragraph()
                a = Tg.part_action_item_dict
                try:
                    action_item_str = a[part][side_correction[side]].get(item)
                except KeyError:
                    raise ActionItemNotFoundError(part, side)
                run = p.add_run()
                if pd.isnull(action_item_str):
                    run.text = 'Keep tracking'
                else:
                    run.text = action_item_str
                    if 'Implementation' in action_item_str:
                        run.font.color.rgb = purple
                    else:
                        run.font.color.rgb = blue
                p.space_before = 0
                p.level = 2

                # p_empty = tf.add_paragraph()
                # p_empty.level = 3
                # p_empty.space_before = 0

    path = Tg.dump_data_root_folder + '\\Monthly_Report\\{}_data\\'.format(show_date.year)
    Tg.chk_dir(path)
    file = 'Conti product monthly yield summary ({}_{}{}).pptx'.format(show_date.year, str(show_date.month).zfill(2),
                                                                       '_' + day if day else '')
    abs_path = path + file
    prs.save(abs_path)
    print('Monthly presentation save to {}'.format(abs_path))


def main():
    # monthly_report_gen(today)
    monthly_report_gen(dt.date(2020, 9, 10), comparing_to_the_last_year=True)


if __name__ == '__main__':
    main()
